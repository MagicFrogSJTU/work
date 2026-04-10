from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image


BASE_DIR = Path(__file__).parent
OUTPUT_FILE = BASE_DIR / "技术分享-20260409-schema17.pptx"
MERMAID_DIR = BASE_DIR / "generated_mermaid"

NAVY = "1F2D3D"
NAVY_ALT = "2B3A4D"
STEEL = "44546A"
ORANGE = "ED7D31"
ORANGE_LIGHT = "F7C9A9"
GOLD = "FFC000"
SAND = "F5F3EF"
FOG = "F3F5F7"
PALE = "ECEFF3"
PALE_BLUE = "E9EEF5"
PALE_ORANGE = "FBE9DE"
INK = "1B1F24"
SLATE = "66788A"
WHITE = "FFFFFF"
SOFT_LINE = "D9DEE5"
MUTED_FILL = "EEF1F4"

TITLE_FONT = "等线 Light"
BODY_FONT = "等线"

SLIDE_W = 13.333
SLIDE_H = 7.5
TOTAL_SLIDES = 17


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape = slide.shapes.add_shape(
        SHAPE.ROUNDED_RECTANGLE if radius else SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x, y, w, h=0.02, color=ORANGE):
    shape = slide.shapes.add_shape(
        SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_text_box(
    slide,
    x,
    y,
    w,
    h,
    paragraphs,
    margin=0.08,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)

    for index, spec in enumerate(paragraphs):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = spec.get("align", PP_ALIGN.LEFT)
        paragraph.space_after = Pt(spec.get("space_after", 3))
        paragraph.space_before = Pt(spec.get("space_before", 0))
        paragraph.line_spacing = spec.get("line_spacing", 1.15)

        run = paragraph.add_run()
        run.text = spec["text"]
        font = run.font
        font.name = spec.get("font", BODY_FONT)
        font.size = Pt(spec.get("size", 18))
        font.bold = spec.get("bold", False)
        font.italic = spec.get("italic", False)
        font.color.rgb = rgb(spec.get("color", INK))

    return box


def add_title(slide, number: int, title: str, dark=False, kicker=None):
    if dark:
        add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
        add_rect(slide, 0, 0, 0.34, SLIDE_H, ORANGE)
        add_rect(slide, 0.82, 0.92, 1.25, 0.34, ORANGE, radius=True)
        add_text_box(
            slide,
            0.92,
            0.97,
            1.05,
            0.2,
            [
                {
                    "text": f"{number:02d}",
                    "size": 11,
                    "bold": True,
                    "color": WHITE,
                    "align": PP_ALIGN.CENTER,
                }
            ],
            margin=0,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if kicker:
            add_text_box(
                slide,
                2.28,
                0.95,
                2.4,
                0.22,
                [
                    {
                        "text": kicker,
                        "size": 11,
                        "bold": True,
                        "color": ORANGE_LIGHT,
                        "align": PP_ALIGN.LEFT,
                    }
                ],
                margin=0,
            )
        add_text_box(
            slide,
            0.92,
            1.45,
            11.2,
            1.45,
            [
                {
                    "text": title,
                    "size": 27,
                    "bold": True,
                    "color": WHITE,
                    "font": TITLE_FONT,
                }
            ],
        )
        add_slide_no(slide, number, light=True)
    else:
        set_bg(slide, WHITE)
        add_rect(slide, 0, 0, SLIDE_W, 0.22, NAVY)
        add_rect(slide, 0.72, 0.55, 0.8, 0.26, ORANGE, radius=True)
        add_text_box(
            slide,
            0.82,
            0.59,
            0.6,
            0.16,
            [
                {
                    "text": f"{number:02d}",
                    "size": 10,
                    "bold": True,
                    "color": WHITE,
                    "align": PP_ALIGN.CENTER,
                }
            ],
            margin=0,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text_box(
            slide,
            1.7,
            0.46,
            9.6,
            0.58,
            [
                {
                    "text": title,
                    "size": 23,
                    "bold": True,
                    "color": INK,
                    "font": TITLE_FONT,
                }
            ],
            margin=0,
        )
        if kicker:
            add_text_box(
                slide,
                1.72,
                0.96,
                2.2,
                0.18,
                [{"text": kicker, "size": 10, "bold": True, "color": ORANGE}],
                margin=0,
            )
        add_line(slide, 0.72, 1.2, 11.0, 0.025, ORANGE)
        add_slide_no(slide, number)


def add_slide_no(slide, number: int, light=False):
    color = WHITE if light else NAVY
    add_text_box(
        slide,
        12.1,
        6.87,
        0.7,
        0.26,
        [
            {
                "text": f"{number:02d}",
                "size": 11,
                "bold": True,
                "color": color,
                "align": PP_ALIGN.RIGHT,
            }
        ],
        margin=0,
    )


def add_chip(slide, x, y, text, fill, color, w=1.45, h=0.34):
    chip = add_rect(slide, x, y, w, h, fill, radius=True)
    add_text_box(
        slide,
        x + 0.05,
        y + 0.01,
        w - 0.1,
        h - 0.02,
        [
            {
                "text": text,
                "size": 11,
                "bold": True,
                "color": color,
                "align": PP_ALIGN.CENTER,
            }
        ],
        margin=0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return chip


def add_stat_box(slide, x, y, w, h, number, label, fill=NAVY, accent=ORANGE):
    add_rect(slide, x, y, w, h, fill, radius=True)
    add_line(slide, x + 0.18, y + 0.18, 0.42, 0.03, accent)
    add_text_box(
        slide,
        x + 0.18,
        y + 0.38,
        w - 0.36,
        0.58,
        [
            {
                "text": number,
                "size": 30,
                "bold": True,
                "color": WHITE,
                "font": TITLE_FONT,
            }
        ],
        margin=0,
    )
    add_text_box(
        slide,
        x + 0.18,
        y + 1.0,
        w - 0.36,
        0.36,
        [{"text": label, "size": 13, "color": "DCE3EA"}],
        margin=0,
    )


def add_card(
    slide,
    x,
    y,
    w,
    h,
    title,
    lines,
    fill=WHITE,
    accent=ORANGE,
    title_color=INK,
    body_color=INK,
    title_size=18,
    body_size=13,
):
    add_rect(slide, x, y, w, h, fill, line=SOFT_LINE, radius=True)
    add_rect(slide, x + 0.14, y + 0.18, 0.09, h - 0.36, accent, radius=True)
    paragraphs = [
        {
            "text": title,
            "size": title_size,
            "bold": True,
            "color": title_color,
            "space_after": 6,
        }
    ]
    for line in lines:
        paragraphs.append(
            {"text": line, "size": body_size, "color": body_color, "space_after": 4}
        )
    add_text_box(slide, x + 0.32, y + 0.16, w - 0.42, h - 0.28, paragraphs)


def add_section_band(slide, x, y, w, h, number, title, subtitle):
    add_rect(slide, x, y, w, h, NAVY, radius=True)
    add_rect(slide, x + 0.22, y + 0.28, 1.05, 0.42, ORANGE, radius=True)
    add_text_box(
        slide,
        x + 0.28,
        y + 0.34,
        0.9,
        0.22,
        [
            {
                "text": number,
                "size": 13,
                "bold": True,
                "color": WHITE,
                "align": PP_ALIGN.CENTER,
            }
        ],
        margin=0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text_box(
        slide,
        x + 0.22,
        y + 0.95,
        w - 0.44,
        0.9,
        [{"text": title, "size": 24, "bold": True, "color": WHITE, "font": TITLE_FONT}],
        margin=0,
    )
    add_text_box(
        slide,
        x + 0.22,
        y + 1.95,
        w - 0.44,
        0.6,
        [{"text": subtitle, "size": 14, "color": "D6DEE7"}],
        margin=0,
    )


def add_dot_bullets(
    slide, x, y, w, items, font_size=15, color=INK, dot=ORANGE, row_gap=0.46
):
    for idx, item in enumerate(items):
        if isinstance(item, dict):
            text = item["text"]
            item_font_size = item.get("size", font_size)
            item_color = item.get("color", color)
            item_dot = item.get("dot", dot)
            indent = item.get("indent", 0)
        else:
            text = item
            item_font_size = font_size
            item_color = color
            item_dot = dot
            indent = 0

        row_y = y + idx * row_gap
        bullet_x = x + indent * 0.34
        bullet_size = 0.1 if indent == 0 else 0.075
        dot_shape = slide.shapes.add_shape(
            SHAPE.OVAL,
            Inches(bullet_x),
            Inches(row_y + 0.1),
            Inches(bullet_size),
            Inches(bullet_size),
        )
        dot_shape.fill.solid()
        dot_shape.fill.fore_color.rgb = rgb(item_dot)
        dot_shape.line.fill.background()
        add_text_box(
            slide,
            bullet_x + 0.18,
            row_y,
            w - (bullet_x - x) - 0.18,
            0.35,
            [{"text": text, "size": item_font_size, "color": item_color}],
            margin=0,
        )


def add_chevron(slide, x, y, w=0.32, h=0.36, fill=ORANGE):
    shape = slide.shapes.add_shape(
        SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()
    return shape


def mermaid_paths(slide_no: int) -> list[Path]:
    return sorted(MERMAID_DIR.glob(f"slide{slide_no:02d}_*.png"))


def fit_image(path: Path, max_w: float, max_h: float) -> tuple[float, float]:
    with Image.open(path) as image:
        width, height = image.size
    ratio = min(max_w / width, max_h / height)
    return width * ratio, height * ratio


def add_image_panel(slide, path: Path, x, y, w, h, caption=None):
    add_rect(slide, x, y, w, h, WHITE, line=SOFT_LINE, radius=True)
    image_w, image_h = fit_image(path, Inches(w - 0.24).emu, Inches(h - 0.56).emu)
    image_w = image_w / Inches(1).emu
    image_h = image_h / Inches(1).emu
    image_x = x + 0.12 + (w - 0.24 - image_w) / 2
    image_y = y + 0.12 + (h - 0.56 - image_h) / 2
    slide.shapes.add_picture(
        str(path),
        Inches(image_x),
        Inches(image_y),
        width=Inches(image_w),
        height=Inches(image_h),
    )
    if caption:
        add_text_box(
            slide,
            x + 0.15,
            y + h - 0.38,
            w - 0.3,
            0.22,
            [
                {
                    "text": caption,
                    "size": 10,
                    "bold": True,
                    "color": SLATE,
                    "align": PP_ALIGN.CENTER,
                }
            ],
            margin=0,
        )


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide, 1, "自动驾驶算法发展切面：\n我的参与与思考", dark=True, kicker="INTRO"
    )
    add_text_box(
        slide,
        0.94,
        2.68,
        5.3,
        0.45,
        [
            {
                "text": "2021-2026  自动驾驶算法 / 数据系统 / 大模型应用",
                "size": 15,
                "color": ORANGE_LIGHT,
            }
        ],
        margin=0,
    )
    add_rect(slide, 0.82, 3.05, 6.15, 2.58, NAVY_ALT, line="41546A", radius=True)
    add_dot_bullets(
        slide,
        1.1,
        3.48,
        5.35,
        [
            "姓名：陈奕志",
            "籍贯：广东",
            "学校：上海交大",
            "工作年限：7年+",
            "上家公司：华为车BU",
        ],
        font_size=17,
        color="E8EEF2",
        dot=ORANGE,
        row_gap=0.48,
    )
    add_stat_box(slide, 8.02, 2.05, 2.0, 1.55, "7Y+", "工作年限")
    add_stat_box(slide, 10.22, 2.05, 2.0, 1.55, "AI", "算法 / 数据 / 系统")
    add_rect(slide, 8.02, 3.92, 4.2, 1.72, WHITE, line="41546A", radius=True)
    add_text_box(
        slide,
        8.28,
        4.18,
        3.2,
        0.34,
        [
            {"text": "分享主线", "size": 14, "bold": True, "color": ORANGE},
        ],
        margin=0,
    )
    add_text_box(
        slide,
        8.28,
        4.68,
        3.58,
        0.72,
        [
            {
                "text": "工作经历回顾 -> 三个方法论心得 -> AI 应用视角总结",
                "size": 13,
                "color": NAVY,
                "line_spacing": 1.05,
            }
        ],
        margin=0,
    )
    return slide


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagram = mermaid_paths(2)[0]
    add_title(slide, 2, "我过去 5 年主要做了什么", kicker="TRACK")
    add_section_band(
        slide,
        0.88,
        1.52,
        3.35,
        4.98,
        "01",
        "工作内容回顾",
        "从感知范式升级，到数据系统建设，再到大模型与 AI 应用。",
    )
    add_image_panel(slide, diagram, 4.58, 1.52, 7.72, 4.98)
    return slide


def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagrams = mermaid_paths(3)
    add_title(slide, 3, "五年，三个心得", kicker="THREE")
    add_card(
        slide,
        0.84,
        1.55,
        3.8,
        4.9,
        "心得 1",
        ["统一模型链路决定上限。"],
        fill=FOG,
        accent=ORANGE,
        title_size=21,
        body_size=12,
    )
    add_image_panel(slide, diagrams[0], 1.08, 3.05, 3.32, 2.48)
    add_card(
        slide,
        4.8,
        1.55,
        3.55,
        4.9,
        "心得 2",
        ["数据定义与闭环决定效率。"],
        fill=FOG,
        accent=STEEL,
        title_size=21,
        body_size=12,
    )
    add_image_panel(slide, diagrams[1], 5.04, 3.05, 3.07, 2.48)
    add_card(
        slide,
        8.52,
        1.55,
        3.95,
        4.9,
        "心得 3",
        ["系统协同决定应用上限。"],
        fill=FOG,
        accent=GOLD,
        title_size=20,
        body_size=12,
    )
    add_image_panel(slide, diagrams[2], 8.82, 3.05, 3.34, 2.48)
    return slide


def build_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagram = mermaid_paths(4)[0]
    add_title(slide, 4, "心得 1：the bitter lesson", kicker="LESSON 1")
    add_rect(slide, 0.82, 1.55, 4.2, 4.95, FOG, line=SOFT_LINE, radius=True)
    add_text_box(
        slide,
        1.12,
        1.86,
        3.55,
        0.42,
        [{"text": "核心判断", "size": 13, "bold": True, "color": ORANGE}],
        margin=0,
    )
    add_dot_bullets(
        slide,
        1.12,
        2.42,
        3.55,
        [
            "算力 + 数据 + 简单通用学习算法",
            "正在持续吃掉人工规则、特征工程和模块接口",
        ],
        font_size=16,
        color=INK,
        dot=ORANGE,
        row_gap=0.86,
    )
    add_rect(slide, 1.12, 4.72, 3.4, 0.9, NAVY, radius=True)
    add_text_box(
        slide,
        1.28,
        4.98,
        3.08,
        0.34,
        [
            {
                "text": "经验判断 -> 趋势判断 -> 组织判断",
                "size": 13,
                "bold": True,
                "color": WHITE,
                "align": PP_ALIGN.CENTER,
            }
        ],
        margin=0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_image_panel(slide, diagram, 5.38, 1.55, 6.82, 4.95)
    return slide


def build_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 5, "心得 1-1：神经网络吃掉传统软件", kicker="2D → 3D")
    add_card(
        slide,
        0.88,
        1.38,
        12.32,
        0.86,
        "一个直接变化",
        [
            "2D 感知时代，核心藏在后处理和融合代码里；3D 感知时代，核心更集中在统一网络与训练数据。"
        ],
        fill=PALE_ORANGE,
        accent=ORANGE,
        title_size=15,
        body_size=14,
    )
    add_image_panel(
        slide, BASE_DIR / "2d_ad_perception.png", 0.85, 2.1, 5.95, 4.7, "2D 感知架构"
    )
    add_image_panel(
        slide, BASE_DIR / "3d_ad_perception.png", 6.9, 2.1, 5.55, 4.7, "3D 感知架构"
    )
    return slide


def build_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagram = mermaid_paths(6)[0]
    add_title(slide, 6, "心得 1-2：越端到端，系统能力上限越高", kicker="E2E")
    add_rect(slide, 0.88, 1.45, 3.55, 5.08, NAVY, radius=True)
    add_text_box(
        slide,
        1.18,
        1.82,
        2.9,
        0.32,
        [{"text": "为什么端到端更强", "size": 14, "bold": True, "color": ORANGE}],
        margin=0,
    )
    add_dot_bullets(
        slide,
        1.18,
        2.38,
        2.85,
        ["链路越长，接口损失和误差传递越重", "端到端，本质是把任务改写成统一学习问题"],
        font_size=15,
        color=WHITE,
        dot=ORANGE,
        row_gap=1.02,
    )
    add_image_panel(slide, diagram, 4.72, 1.72, 7.5, 4.42)
    return slide


def build_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagrams = mermaid_paths(7)
    add_title(slide, 7, "心得 1-3：简单留给网络，复杂留给数据", kicker="DATA OPS")
    add_card(
        slide,
        0.9,
        1.45,
        4.75,
        2.35,
        "模型侧",
        ["架构更统一、更简洁", "文本 / 语音 / 视觉，分类 / 回归 / 生成持续收敛"],
        fill=FOG,
        accent=STEEL,
        title_size=18,
        body_size=14,
    )
    add_card(
        slide,
        0.9,
        4.0,
        4.75,
        2.15,
        "数据侧",
        ["数据规模、团队规模、研发复杂度持续上升", "真正难点转移到数据工程与产线管理"],
        fill=PALE_ORANGE,
        accent=ORANGE,
        title_size=18,
        body_size=14,
    )
    add_image_panel(slide, diagrams[0], 5.95, 1.45, 2.7, 1.55)
    add_rect(slide, 5.95, 3.18, 2.7, 2.95, PALE_BLUE, line=SOFT_LINE, radius=True)
    add_text_box(
        slide,
        6.18,
        3.45,
        2.2,
        0.26,
        [{"text": "一个例子", "size": 15, "bold": True, "color": STEEL}],
        margin=0,
    )
    add_text_box(
        slide,
        6.18,
        3.9,
        2.18,
        1.45,
        [
            {
                "text": "静态感知里，在线推理变得更简，离线数据重建与标注链路反而更长。",
                "size": 13,
                "color": INK,
            }
        ],
        margin=0,
    )
    add_image_panel(slide, diagrams[1], 8.88, 1.45, 3.6, 4.95)
    return slide


def build_slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagram = mermaid_paths(8)[0]
    add_title(slide, 8, "心得 2：当我们说“数据”时，是在说什么", kicker="LESSON 2")
    add_section_band(
        slide,
        0.92,
        1.62,
        3.25,
        4.82,
        "02",
        "数据到底是什么",
        "数据不是素材集合，而是任务定义、闭环流程和可计算资产的总和。",
    )
    add_image_panel(slide, diagram, 4.52, 1.62, 7.72, 4.82)
    return slide


def build_slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagram = mermaid_paths(9)[0]
    add_title(slide, 9, "心得 2-1：数据定义，是算法问题", kicker="TASK")
    add_card(
        slide,
        0.9,
        1.5,
        5.1,
        4.95,
        "结论",
        [
            "输入输出，定义了“你要达到什么目的”，即任务本身。",
            "评测指标，定义了“你认为什么是好的”，即任务的优化目标。",
            "抛开这三个谈算法，是无源之水，镜中之月。",
        ],
        fill=FOG,
        accent=ORANGE,
        title_size=20,
        body_size=15,
    )
    add_image_panel(slide, diagram, 6.3, 1.78, 5.92, 4.15)
    return slide


def build_slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagram = mermaid_paths(10)[0]
    add_title(slide, 10, "心得 2-2：数据闭环，是工程问题", kicker="LOOP")
    add_rect(slide, 0.82, 1.45, 4.9, 5.25, FOG, line=SOFT_LINE, radius=True)
    add_dot_bullets(
        slide,
        1.06,
        1.82,
        4.25,
        [
            {"text": "数据理解：标注、统计、可视化、结构化", "size": 14},
            {"text": "数据策划：挖掘、采集、合成、平衡", "size": 14},
            {"text": "人工标注：", "size": 14, "dot": ORANGE},
            {
                "text": "质量重于数量. scaleAI/surgeAI标注员时薪可至200美元.",
                "size": 12,
                "indent": 1,
            },
            {"text": "成本高昂", "size": 12, "indent": 1},
            {
                "text": "基于不信任的团队管理：规格/产线管理/工艺/质量",
                "size": 11,
                "indent": 2,
            },
            {
                "text": "基础设施：标注平台，人机共标算法，统计/运维后台",
                "size": 11,
                "indent": 2,
            },
        ],
        font_size=14,
        row_gap=0.58,
    )
    add_image_panel(slide, diagram, 5.95, 1.58, 6.25, 4.55)
    return slide


def build_slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagram = mermaid_paths(11)[0]
    add_title(slide, 11, "心得 2-3：训练数据没有量，就没有真正的价值", kicker="SCALE")
    add_stat_box(slide, 0.95, 1.55, 2.28, 1.5, "TB/PB", "数据规模")
    add_stat_box(slide, 3.48, 1.55, 2.28, 1.5, "Version", "版本资产")
    add_rect(slide, 0.95, 3.35, 4.8, 2.95, FOG, line=SOFT_LINE, radius=True)
    add_dot_bullets(
        slide,
        1.18,
        3.72,
        4.3,
        [
            "训练数据价值，取决于存储、版本、检索、计算",
            "大规模数据资产，决定模型迭代速度",
            "我做过版本系统、检索系统、多模态数据仓和自动化产线",
        ],
        font_size=14,
        row_gap=0.68,
    )
    add_image_panel(slide, diagram, 6.1, 1.55, 6.12, 4.95)
    return slide


def build_slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagram = mermaid_paths(12)[0]
    add_title(slide, 12, "心得 3：AI 应用是系统性工程", kicker="LESSON 3")
    add_section_band(
        slide,
        0.9,
        1.58,
        3.4,
        4.82,
        "03",
        "系统协同",
        "AI 应用的核心不是单点最强，而是系统协同最优；大模型时代，这进一步体现为知识瓶颈和上下文瓶颈。",
    )
    add_image_panel(slide, diagram, 4.72, 1.58, 7.5, 4.82)
    return slide


def build_slide_13(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 13, "心得 3-1：取长补短，系统协同", kicker="SYSTEM")
    add_card(
        slide,
        0.92,
        1.52,
        5.0,
        5.12,
        "系统协同的含义",
        [
            "超人模块 + 弱智 AI，也能做出可用系统",
            "激光雷达和 RoadCode（华为众包高精地图），都是典型例子",
            "系统的竞争力，本质来自协同、互助。",
        ],
        fill=FOG,
        accent=ORANGE,
        title_size=20,
        body_size=15,
    )
    add_image_panel(
        slide,
        BASE_DIR / "autonomous_driving_arch.png",
        6.15,
        1.55,
        6.15,
        5.15,
        "自动驾驶系统架构",
    )
    return slide


def build_slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagrams = mermaid_paths(14)
    add_title(slide, 14, "心得 3-2：预训练是大模型的知识瓶颈", kicker="PRETRAIN")
    add_card(
        slide,
        0.9,
        1.35,
        12.32,
        1.0,
        "判断",
        [
            "冷门垂域的微调有用，但瓶颈仍然在预训练；真正的竞争力来自能否把垂域知识持续送回更大规模训练。"
        ],
        fill=PALE_ORANGE,
        accent=ORANGE,
        title_size=15,
        body_size=14,
    )
    add_image_panel(slide, diagrams[0], 1.0, 2.55, 5.4, 3.6)
    add_image_panel(slide, diagrams[1], 6.85, 2.55, 5.4, 3.6)
    return slide


def build_slide_15(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagram = mermaid_paths(15)[0]
    add_title(slide, 15, "心得 3-3：上下文窗口是大模型的应用瓶颈", kicker="CONTEXT")
    add_rect(slide, 0.92, 1.55, 4.62, 4.92, NAVY, radius=True)
    add_text_box(
        slide,
        1.22,
        1.95,
        3.8,
        0.34,
        [{"text": "应用层矛盾", "size": 14, "bold": True, "color": ORANGE}],
        margin=0,
    )
    add_dot_bullets(
        slide,
        1.22,
        2.45,
        3.65,
        [
            "大模型没有记忆，只有上下文",
            "上下文工程的本质，在减少单个 agent 的上下文负担",
        ],
        font_size=15,
        color=WHITE,
        dot=ORANGE,
        row_gap=1.0,
    )
    add_image_panel(slide, diagram, 5.78, 1.55, 6.42, 4.92)
    return slide


def build_slide_16(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    diagrams = mermaid_paths(16)
    add_title(slide, 16, "总结", kicker="WRAP-UP")
    add_stat_box(slide, 0.95, 1.45, 2.2, 1.35, "01", "数据是核心竞争力")
    add_stat_box(slide, 3.4, 1.45, 2.2, 1.35, "02", "垂域壁垒会长期存在", fill=STEEL)
    add_image_panel(slide, diagrams[0], 0.9, 2.95, 6.1, 3.62)
    add_image_panel(slide, diagrams[1], 7.08, 2.95, 5.32, 3.62)
    return slide


def build_slide_17(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 17, "展望未来", dark=True, kicker="FUTURE")
    add_rect(slide, 0.82, 2.15, 6.9, 3.12, NAVY_ALT, line="41546A", radius=True)
    add_dot_bullets(
        slide,
        1.12,
        2.62,
        6.15,
        [
            {
                "text": "中远海运 x AI，机会巨大，值得期待！",
                "size": 21,
                "color": WHITE,
                "dot": ORANGE,
            },
            {"text": "预祝研究院做大做强！", "size": 21, "color": WHITE, "dot": ORANGE},
            {"text": "ANY QUESTIONS?", "size": 24, "color": GOLD, "dot": GOLD},
        ],
        font_size=21,
        color=WHITE,
        dot=ORANGE,
        row_gap=0.8,
    )
    add_stat_box(
        slide, 8.35, 2.3, 3.1, 1.55, "SHIP + AI", "行业深度 x 数据体系 x 系统工程"
    )
    add_rect(slide, 8.35, 4.2, 3.1, 1.18, WHITE, line="41546A", radius=True)
    add_text_box(
        slide,
        8.58,
        4.54,
        2.62,
        0.34,
        [
            {
                "text": "下一阶段真正稀缺的是可持续落地能力。",
                "size": 14,
                "bold": True,
                "color": NAVY,
                "align": PP_ALIGN.CENTER,
            }
        ],
        margin=0,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return slide


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.author = "GitHub Copilot"
    prs.core_properties.title = "自动驾驶工作复盘分享"
    prs.core_properties.subject = "华为自动驾驶经历与 AI 心得"

    builders = [
        build_slide_1,
        build_slide_2,
        build_slide_3,
        build_slide_4,
        build_slide_5,
        build_slide_6,
        build_slide_7,
        build_slide_8,
        build_slide_9,
        build_slide_10,
        build_slide_11,
        build_slide_12,
        build_slide_13,
        build_slide_14,
        build_slide_15,
        build_slide_16,
        build_slide_17,
    ]
    for builder in builders:
        builder(prs)
    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(str(OUTPUT_FILE))
    print(f"Wrote {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
